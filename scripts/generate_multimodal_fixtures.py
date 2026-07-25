"""One-off generator for the multimodal extraction test fixtures.

Not part of the test suite — run manually whenever the fixtures need
regenerating:

    python scripts/generate_multimodal_fixtures.py

Produces, under tests/fixtures/:
    multimodal_sample.pdf   (text + table + image)
    multimodal_sample.docx  (text + table + image)
    multimodal_sample.pptx  (text + table + image)
    multimodal_sample.md    (text + table; Markdown has no embeddable
                             binary image, so no image element is expected
                             for this fixture)
"""

import io
import os

from PIL import Image
from docx import Document as DocxDocument
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.units import inch
from reportlab.platypus import (
    SimpleDocTemplate,
    Paragraph,
    Table,
    TableStyle,
    Image as RLImage,
    Spacer,
)
from reportlab.lib.styles import getSampleStyleSheet

FIXTURES_DIR = os.path.join(os.path.dirname(__file__), "..", "tests", "fixtures")


def _sample_image_path(tmp_dir: str) -> str:
    """Create a small red/blue PNG on disk and return its path."""
    path = os.path.join(tmp_dir, "_sample_image.png")
    image = Image.new("RGB", (64, 64), color="red")
    for x in range(32, 64):
        for y in range(64):
            image.putpixel((x, y), (0, 0, 255))
    image.save(path, format="PNG")
    return path


def generate_pdf(image_path: str) -> None:
    out_path = os.path.join(FIXTURES_DIR, "multimodal_sample.pdf")
    styles = getSampleStyleSheet()
    doc = SimpleDocTemplate(out_path, pagesize=letter)

    table_data = [["Name", "Score"], ["Alice", "92"], ["Bob", "85"]]
    table = Table(table_data, colWidths=[2 * inch, 2 * inch])
    table.setStyle(
        TableStyle(
            [
                ("GRID", (0, 0), (-1, -1), 1, colors.black),
                ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
                ("FONTSIZE", (0, 0), (-1, -1), 12),
                ("TOPPADDING", (0, 0), (-1, -1), 8),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
            ]
        )
    )

    story = [
        Paragraph("Multimodal Sample Document", styles["Title"]),
        Spacer(1, 0.4 * inch),
        Paragraph(
            "This is a sample paragraph of narrative text used to exercise "
            "text extraction in the multimodal document embedding tool.",
            styles["BodyText"],
        ),
        Spacer(1, 0.6 * inch),
        table,
        Spacer(1, 0.6 * inch),
        RLImage(image_path, width=100, height=100),
    ]
    doc.build(story)


def generate_docx(image_path: str) -> None:
    out_path = os.path.join(FIXTURES_DIR, "multimodal_sample.docx")
    document = DocxDocument()
    document.add_heading("Multimodal Sample Document", level=1)
    document.add_paragraph(
        "This is a sample paragraph of narrative text used to exercise "
        "text extraction in the multimodal document embedding tool."
    )

    table = document.add_table(rows=3, cols=2)
    values = [("Name", "Score"), ("Alice", "92"), ("Bob", "85")]
    for row, (name, score) in zip(table.rows, values):
        row.cells[0].text = name
        row.cells[1].text = score

    document.add_picture(image_path, width=None)
    document.save(out_path)


def generate_pptx(image_path: str) -> None:
    out_path = os.path.join(FIXTURES_DIR, "multimodal_sample.pptx")
    presentation = Presentation()

    text_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    text_slide.shapes.title.text = "Multimodal Sample Document"
    text_slide.placeholders[1].text = (
        "This is a sample paragraph of narrative text used to exercise "
        "text extraction in the multimodal document embedding tool."
    )

    table_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    rows, cols = 3, 2
    table_shape = table_slide.shapes.add_table(
        rows, cols, Inches(1), Inches(1.5), Inches(4), Inches(2)
    )
    table = table_shape.table
    values = [("Name", "Score"), ("Alice", "92"), ("Bob", "85")]
    for row_index, (name, score) in enumerate(values):
        table.cell(row_index, 0).text = name
        table.cell(row_index, 1).text = score

    image_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    image_slide.shapes.add_picture(image_path, Inches(1), Inches(1), height=Inches(2))

    presentation.save(out_path)


def generate_md() -> None:
    out_path = os.path.join(FIXTURES_DIR, "multimodal_sample.md")
    content = """# Multimodal Sample Document

This is a sample paragraph of narrative text used to exercise text
extraction in the multimodal document embedding tool.

| Name  | Score |
|-------|-------|
| Alice | 92    |
| Bob   | 85    |
"""
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(content)


def main() -> None:
    os.makedirs(FIXTURES_DIR, exist_ok=True)
    image_path = _sample_image_path(FIXTURES_DIR)
    try:
        generate_pdf(image_path)
        generate_docx(image_path)
        generate_pptx(image_path)
        generate_md()
    finally:
        if os.path.exists(image_path):
            os.remove(image_path)


if __name__ == "__main__":
    main()
